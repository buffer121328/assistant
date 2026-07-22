from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from email.message import EmailMessage
import os
from pathlib import Path
import re
import tempfile
from typing import Callable, TypeVar

_SAFE_TASK_ID = re.compile(r"^[A-Za-z0-9][A-Za-z0-9_-]{0,127}$")
_SAFE_FILENAME = re.compile(r"^[^/\\\x00]{1,128}$")
_T = TypeVar("_T")


class ArtifactPathError(ValueError):
    """表示 处理 artifact path error 的后端数据结构或服务对象。"""

    pass


class OptionalOfficeDependencyError(RuntimeError):
    """表示 处理 optional office dependency error 的后端数据结构或服务对象。"""

    def __init__(self) -> None:
        """初始化对象实例。"""
        super().__init__(
            "Optional Office dependencies are not installed. "
            "Install with: uv sync --extra office"
        )


@dataclass(frozen=True)
class Artifact:
    """表示 处理 artifact 的后端数据结构或服务对象。"""

    reference: str
    media_type: str
    size_bytes: int


class ArtifactStore:
    """表示 处理 artifact store 的后端数据结构或服务对象。"""

    def __init__(self, root: Path) -> None:
        """初始化对象实例。

        Args:
            root: root 参数。
        """
        self.root = root.expanduser().resolve()
        self.root.mkdir(parents=True, exist_ok=True)

    def reserve(self, *, task_id: str, filename: str, suffix: str) -> Path:
        """处理 reserve。

        Args:
            task_id: task_id 参数。
            filename: filename 参数。
            suffix: suffix 参数。
        """
        safe_task_id = task_id.strip()
        safe_filename = filename.strip()
        if not _SAFE_TASK_ID.fullmatch(safe_task_id):
            raise ArtifactPathError("Invalid task id")
        if (
            not _SAFE_FILENAME.fullmatch(safe_filename)
            or Path(safe_filename).is_absolute()
            or Path(safe_filename).name != safe_filename
            or safe_filename in {".", ".."}
        ):
            raise ArtifactPathError("Invalid artifact filename")
        if (
            not suffix.startswith(".")
            or Path(safe_filename).suffix.lower() != suffix.lower()
        ):
            raise ArtifactPathError(f"Artifact filename must end with {suffix}")

        task_root = self.root / safe_task_id
        if task_root.is_symlink():
            raise ArtifactPathError("Task artifact directory must not be a symlink")
        task_root.mkdir(mode=0o700, parents=False, exist_ok=True)
        resolved_task_root = task_root.resolve(strict=True)
        if resolved_task_root.parent != self.root:
            raise ArtifactPathError("Task artifact directory escaped root")

        target = resolved_task_root / safe_filename
        if target.is_symlink() or target.parent != resolved_task_root:
            raise ArtifactPathError("Artifact path escaped task directory")
        return target

    def atomic_write_bytes(self, target: Path, data: bytes) -> None:
        """处理 atomic write bytes。

        Args:
            target: target 参数。
            data: data 参数。
        """
        self._assert_target(target)
        descriptor, temporary_name = tempfile.mkstemp(
            dir=target.parent,
            prefix=f".{target.name}.",
            suffix=".tmp",
        )
        temporary = Path(temporary_name)
        try:
            with os.fdopen(descriptor, "wb") as stream:
                stream.write(data)
                stream.flush()
                os.fsync(stream.fileno())
            os.replace(temporary, target)
        finally:
            temporary.unlink(missing_ok=True)

    def atomic_save(self, target: Path, save: Callable[[str], _T]) -> _T:
        """处理 atomic save。

        Args:
            target: target 参数。
            save: save 参数。
        """
        self._assert_target(target)
        descriptor, temporary_name = tempfile.mkstemp(
            dir=target.parent,
            prefix=f".{target.name}.",
            suffix=target.suffix,
        )
        os.close(descriptor)
        temporary = Path(temporary_name)
        try:
            result = save(str(temporary))
            os.replace(temporary, target)
            return result
        finally:
            temporary.unlink(missing_ok=True)

    def describe(self, target: Path, *, media_type: str) -> Artifact:
        """处理 describe。

        Args:
            target: target 参数。
            media_type: media_type 参数。
        """
        self._assert_target(target)
        relative = target.relative_to(self.root).as_posix()
        return Artifact(
            reference=relative,
            media_type=media_type,
            size_bytes=target.stat().st_size,
        )

    def absolute_path(self, *, task_id: str, reference: str) -> Path:
        """处理 absolute path。

        Args:
            task_id: task_id 参数。
            reference: reference 参数。
        """
        if not _SAFE_TASK_ID.fullmatch(task_id.strip()):
            raise ArtifactPathError("Invalid task id")
        candidate = self.root / reference
        resolved = candidate.resolve(strict=True)
        task_root = (self.root / task_id.strip()).resolve(strict=True)
        if not resolved.is_relative_to(task_root) or resolved.parent != task_root:
            raise ArtifactPathError("Artifact reference escaped task directory")
        return resolved

    def read_bytes(self, *, task_id: str, reference: str) -> bytes:
        """处理 read bytes。

        Args:
            task_id: task_id 参数。
            reference: reference 参数。
        """
        return self.absolute_path(task_id=task_id, reference=reference).read_bytes()

    def read_text(self, *, task_id: str, reference: str) -> str:
        """处理 read text。

        Args:
            task_id: task_id 参数。
            reference: reference 参数。
        """
        return self.absolute_path(task_id=task_id, reference=reference).read_text(
            encoding="utf-8"
        )

    def _assert_target(self, target: Path) -> None:
        """执行 处理 assert target 的内部辅助逻辑。

        Args:
            target: target 参数。
        """
        resolved_parent = target.parent.resolve(strict=True)
        if resolved_parent.parent != self.root or target.name != target.name.strip():
            raise ArtifactPathError("Artifact target escaped task directory")
        if target.is_symlink():
            raise ArtifactPathError("Artifact target must not be a symlink")


class ProductivityTools:
    """表示 处理 productivity tools 的后端数据结构或服务对象。"""

    def __init__(self, store: ArtifactStore) -> None:
        """初始化对象实例。

        Args:
            store: store 参数。
        """
        self.store = store

    def create_email_draft(
        self,
        *,
        task_id: str,
        filename: str,
        subject: str,
        body: str,
        to: tuple[str, ...] = (),
    ) -> Artifact:
        """创建 email draft。

        Args:
            task_id: task_id 参数。
            filename: filename 参数。
            subject: subject 参数。
            body: body 参数。
            to: to 参数。
        """
        message = EmailMessage()
        message["Subject"] = _bounded_text(subject, "subject", 300)
        if to:
            message["To"] = ", ".join(
                _bounded_text(item, "recipient", 320) for item in to[:20]
            )
        message.set_content(_bounded_text(body, "body", 100_000))
        target = self.store.reserve(task_id=task_id, filename=filename, suffix=".eml")
        self.store.atomic_write_bytes(target, message.as_bytes())
        return self.store.describe(target, media_type="message/rfc822")

    def create_calendar_event(
        self,
        *,
        task_id: str,
        filename: str,
        title: str,
        start: str,
        end: str,
        description: str = "",
    ) -> Artifact:
        """创建 calendar event。

        Args:
            task_id: task_id 参数。
            filename: filename 参数。
            title: title 参数。
            start: start 参数。
            end: end 参数。
            description: description 参数。
        """
        start_at = datetime.fromisoformat(start)
        end_at = datetime.fromisoformat(end)
        if end_at <= start_at:
            raise ValueError("Calendar event end must be after start")
        content = "\r\n".join(
            (
                "BEGIN:VCALENDAR",
                "VERSION:2.0",
                "PRODID:-//Personal Agent//EN",
                "BEGIN:VEVENT",
                f"SUMMARY:{_ics_escape(_bounded_text(title, 'title', 500))}",
                f"DTSTART:{_ics_datetime(start_at)}",
                f"DTEND:{_ics_datetime(end_at)}",
                f"DESCRIPTION:{_ics_escape(_optional_bounded_text(description, 'description', 10_000))}",
                "END:VEVENT",
                "END:VCALENDAR",
                "",
            )
        )
        target = self.store.reserve(task_id=task_id, filename=filename, suffix=".ics")
        self.store.atomic_write_bytes(target, content.encode("utf-8"))
        return self.store.describe(target, media_type="text/calendar")

    def create_docx(
        self,
        *,
        task_id: str,
        filename: str,
        title: str,
        paragraphs: tuple[str, ...],
    ) -> Artifact:
        """创建 docx。

        Args:
            task_id: task_id 参数。
            filename: filename 参数。
            title: title 参数。
            paragraphs: paragraphs 参数。
        """
        try:
            from docx import Document
        except ImportError as exc:
            raise OptionalOfficeDependencyError() from exc

        document = Document()
        document.add_heading(_bounded_text(title, "title", 500), level=0)
        for paragraph in paragraphs[:100]:
            document.add_paragraph(_bounded_text(paragraph, "paragraph", 10_000))
        target = self.store.reserve(task_id=task_id, filename=filename, suffix=".docx")
        self.store.atomic_save(target, document.save)
        return self.store.describe(
            target,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

    def create_xlsx(
        self,
        *,
        task_id: str,
        filename: str,
        sheet_name: str,
        rows: tuple[tuple[object, ...], ...],
    ) -> Artifact:
        """创建 xlsx。

        Args:
            task_id: task_id 参数。
            filename: filename 参数。
            sheet_name: sheet_name 参数。
            rows: rows 参数。
        """
        try:
            from openpyxl import Workbook  # type: ignore[import-untyped]
        except ImportError as exc:
            raise OptionalOfficeDependencyError() from exc

        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = _bounded_text(sheet_name, "sheet name", 31)
        for row in rows[:10_000]:
            worksheet.append(tuple(_cell_value(item) for item in row[:100]))
        target = self.store.reserve(task_id=task_id, filename=filename, suffix=".xlsx")
        self.store.atomic_save(target, workbook.save)
        return self.store.describe(
            target,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    def create_pptx(
        self,
        *,
        task_id: str,
        filename: str,
        title: str,
        slides: tuple[tuple[str, tuple[str, ...]], ...],
    ) -> Artifact:
        """创建 pptx。

        Args:
            task_id: task_id 参数。
            filename: filename 参数。
            title: title 参数。
            slides: slides 参数。
        """
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise OptionalOfficeDependencyError() from exc

        presentation = Presentation()
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = _bounded_text(title, "title", 500)
        for heading, bullets in slides[:50]:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = _bounded_text(heading, "heading", 500)
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for index, bullet in enumerate(bullets[:50]):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = _bounded_text(bullet, "bullet", 2_000)
        target = self.store.reserve(task_id=task_id, filename=filename, suffix=".pptx")
        self.store.atomic_save(target, presentation.save)
        return self.store.describe(
            target,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )


def _bounded_text(value: object, name: str, maximum: int) -> str:
    """执行 处理 bounded text 的内部辅助逻辑。

    Args:
        value: value 参数。
        name: name 参数。
        maximum: maximum 参数。
    """
    if not isinstance(value, str):
        raise ValueError(f"{name} must be text")
    normalized = value.strip()
    if not normalized or len(normalized) > maximum:
        raise ValueError(f"{name} is empty or too long")
    return normalized


def _optional_bounded_text(value: object, name: str, maximum: int) -> str:
    """执行 处理 optional bounded text 的内部辅助逻辑。

    Args:
        value: value 参数。
        name: name 参数。
        maximum: maximum 参数。
    """
    if not isinstance(value, str):
        raise ValueError(f"{name} must be text")
    normalized = value.strip()
    if len(normalized) > maximum:
        raise ValueError(f"{name} is too long")
    return normalized


def _cell_value(value: object) -> str | int | float | bool | None:
    """执行 处理 cell value 的内部辅助逻辑。

    Args:
        value: value 参数。
    """
    if value is None or isinstance(value, str | int | float | bool):
        return value
    return str(value)[:10_000]


def _ics_escape(value: str) -> str:
    """执行 处理 ics escape 的内部辅助逻辑。

    Args:
        value: value 参数。
    """
    return (
        value.replace("\\", "\\\\")
        .replace(";", "\\;")
        .replace(",", "\\,")
        .replace("\n", "\\n")
    )


def _ics_datetime(value: datetime) -> str:
    """执行 处理 ics datetime 的内部辅助逻辑。

    Args:
        value: value 参数。
    """
    if value.tzinfo is not None:
        return value.astimezone().strftime("%Y%m%dT%H%M%SZ")
    return value.strftime("%Y%m%dT%H%M%S")
