from __future__ import annotations

from email import policy
from email.parser import BytesParser
from pathlib import Path

from docx import Document
from openpyxl import load_workbook  # type: ignore[import-untyped]
from pptx import Presentation
import pytest

from packages.tools.artifacts import ArtifactPathError, ArtifactStore, ProductivityTools
from packages.tools.browser import BrowserDestinationError, PublicUrlPolicy
from packages.tools.sandbox import DockerSandboxConfig, DockerSandboxRunner
from packages.tools.personal import build_personal_tool_descriptors


def test_personal_productivity_tools_generate_real_task_scoped_files(
    tmp_path: Path,
) -> None:
    store = ArtifactStore(tmp_path)
    tools = ProductivityTools(store)

    eml = tools.create_email_draft(
        task_id="task-1",
        filename="draft.eml",
        subject="周报",
        body="本周完成。",
        to=("owner@example.com",),
    )
    ics = tools.create_calendar_event(
        task_id="task-1",
        filename="meeting.ics",
        title="项目复盘",
        start="2026-07-15T09:00:00+08:00",
        end="2026-07-15T10:00:00+08:00",
    )
    docx = tools.create_docx(
        task_id="task-1",
        filename="report.docx",
        title="项目报告",
        paragraphs=("第一段", "第二段"),
    )
    xlsx = tools.create_xlsx(
        task_id="task-1",
        filename="data.xlsx",
        sheet_name="数据",
        rows=(("事项", "状态"), ("测试", "完成")),
    )
    pptx = tools.create_pptx(
        task_id="task-1",
        filename="slides.pptx",
        title="进展",
        slides=(("本周", ("完成工具层", "补充测试")),),
    )

    task_root = (tmp_path / "task-1").resolve()
    for artifact in (eml, ics, docx, xlsx, pptx):
        path = store.absolute_path(task_id="task-1", reference=artifact.reference)
        assert path.is_relative_to(task_root)
        assert path.exists()

    message = BytesParser(policy=policy.default).parsebytes(
        store.read_bytes(task_id="task-1", reference=eml.reference)
    )
    assert message["Subject"] == "周报"
    assert "BEGIN:VEVENT" in store.read_text(task_id="task-1", reference=ics.reference)
    assert Document(str(store.absolute_path(task_id="task-1", reference=docx.reference))).paragraphs[0].text == "项目报告"
    assert load_workbook(store.absolute_path(task_id="task-1", reference=xlsx.reference))["数据"]["A2"].value == "测试"
    assert Presentation(str(store.absolute_path(task_id="task-1", reference=pptx.reference))).slides[0].shapes.title.text == "进展"


@pytest.mark.parametrize("filename", ["../escape.docx", "/tmp/escape.docx", "other/task.docx"])
def test_artifact_store_rejects_path_escape(tmp_path: Path, filename: str) -> None:
    store = ArtifactStore(tmp_path)

    with pytest.raises(ArtifactPathError):
        store.reserve(task_id="task-1", filename=filename, suffix=".docx")

    assert not (tmp_path.parent / "escape.docx").exists()


@pytest.mark.asyncio
async def test_public_url_policy_blocks_private_and_accepts_public() -> None:
    async def resolver(host: str) -> tuple[str, ...]:
        return {
            "public.example": ("93.184.216.34",),
            "private.example": ("127.0.0.1",),
        }[host]

    policy = PublicUrlPolicy(resolver=resolver)

    assert await policy.validate("https://public.example/page") == "https://public.example/page"
    with pytest.raises(BrowserDestinationError):
        await policy.validate("http://private.example/internal")
    with pytest.raises(BrowserDestinationError):
        await policy.validate("file:///etc/passwd")


def test_docker_sandbox_uses_fixed_isolation_flags_and_no_host_fallback(
    tmp_path: Path,
) -> None:
    config = DockerSandboxConfig(
        enabled=True,
        image="python:3.12-alpine",
        allowed_images=("python:3.12-alpine",),
        memory_mb=128,
        cpu_count=0.5,
        pids_limit=32,
    )
    runner = DockerSandboxRunner(config=config, workspace_root=tmp_path)

    argv = runner.build_argv(task_id="task-1", command=("python", "-V"))

    assert argv[:2] == ("docker", "run")
    assert "--network" in argv and "none" in argv
    assert "--read-only" in argv
    assert ("--cap-drop", "ALL") == (argv[argv.index("--cap-drop")], argv[argv.index("--cap-drop") + 1])
    assert "no-new-privileges" in argv
    assert str((tmp_path / "task-1").resolve()) in " ".join(argv)
    assert "/var/run/docker.sock" not in " ".join(argv)

    disabled = DockerSandboxRunner(
        config=DockerSandboxConfig(enabled=False),
        workspace_root=tmp_path,
    )
    assert disabled.available is False
    with pytest.raises(RuntimeError, match="disabled"):
        disabled.build_argv(task_id="task-1", command=("uname", "-a"))


def test_optional_external_tools_are_unavailable_by_default() -> None:
    descriptors = build_personal_tool_descriptors(
        browser_available=False,
        sandbox_available=False,
    )
    by_name = {item.name: item for item in descriptors}

    assert by_name["email.draft"].enabled is True
    assert by_name["office.create_docx"].parallel_safe is True
    assert by_name["browser.read"].enabled is False
    assert by_name["shell.exec"].enabled is False
    assert by_name["shell.exec"].risk_level == "L3"
    assert "email.send" not in by_name
    assert "calendar.sync_event" not in by_name
