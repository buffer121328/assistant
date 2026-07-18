from __future__ import annotations

from pathlib import Path


PARSER_VERSION = "extract-v1"
SUPPORTED_MEDIA_TYPES = {
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


class ExtractionError(RuntimeError):
    pass


class OptionalOfficeDependencyError(ExtractionError):
    pass


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    try:
        if suffix in {".txt", ".md"}:
            text = path.read_text(encoding="utf-8")
        elif suffix == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError as exc:
                raise OptionalOfficeDependencyError("knowledge_optional_office_missing") from exc

            text = "\n".join(page.extract_text() or "" for page in PdfReader(str(path)).pages)
        elif suffix == ".docx":
            try:
                from docx import Document
            except ImportError as exc:
                raise OptionalOfficeDependencyError("knowledge_optional_office_missing") from exc

            text = "\n".join(paragraph.text for paragraph in Document(str(path)).paragraphs)
        elif suffix == ".xlsx":
            text = _xlsx_text(path)
        elif suffix == ".pptx":
            text = _pptx_text(path)
        else:
            raise ExtractionError("knowledge_type_unsupported")
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError("knowledge_parse_failed") from exc
    normalized = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not normalized:
        raise ExtractionError("knowledge_content_empty")
    return normalized


def _xlsx_text(path: Path) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError as exc:
        raise OptionalOfficeDependencyError("knowledge_optional_office_missing") from exc

    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        lines: list[str] = []
        for worksheet in workbook.worksheets:
            lines.append(worksheet.title)
            for row in worksheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value is not None]
                if values:
                    lines.append("\t".join(values))
        return "\n".join(lines)
    finally:
        workbook.close()


def _pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise OptionalOfficeDependencyError("knowledge_optional_office_missing") from exc

    lines: list[str] = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                lines.append(str(text))
    return "\n".join(lines)
